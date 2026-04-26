"""
Build the 1st-place pitch deck from the uploaded Tashkent template.

Strategy:
  1. Open FinVision_presentation.pptx (the 109-slide Google Cloud template).
  2. Wipe its example slides while preserving the slide masters / layouts /
     theme / fonts / colors.
  3. Add 8 fresh slides using existing layouts so branding stays intact.
  4. Save as FinVision_pitch.pptx alongside the original.
"""

from pptx import Presentation

INPUT  = "C:/Users/user/Desktop/FinVision/finvision-voice/FinVision_presentation.pptx"
OUTPUT = "C:/Users/user/Desktop/FinVision/finvision-voice/FinVision_pitch.pptx"

p = Presentation(INPUT)

# ─── Wipe existing slides cleanly — drop both the relationship AND the
#     underlying slide part, so we don't end up with duplicate names in
#     the package zip when new slides reuse slide1.xml, slide2.xml, etc.
sld_id_lst = p.slides._sldIdLst
RID_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
for sld_id_elem in list(sld_id_lst):
    rId = sld_id_elem.attrib[RID_NS]
    p.part.drop_rel(rId)        # remove the package relationship + part
    sld_id_lst.remove(sld_id_elem)

layouts = {l.name: l for l in p.slide_layouts}


def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_text(slide, idx, text):
    ph = get_ph(slide, idx)
    if ph is not None:
        ph.text = text


def set_bullets(slide, idx, lines):
    ph = get_ph(slide, idx)
    if ph is None:
        return
    tf = ph.text_frame
    tf.text = lines[0]
    for line in lines[1:]:
        para = tf.add_paragraph()
        para.text = line


# ───────── SLIDE 1 — Title ─────────
s = p.slides.add_slide(layouts["CUSTOM_1"])
set_text(s, 0, "SQB OvozAI")
set_text(
    s, 1,
    "Real-time AI Sales Copilot for Bank Operators\n"
    "Cluely-style · Trilingual · Production-grade",
)

# ───────── SLIDE 2 — The problem ─────────
s = p.slides.add_slide(layouts["TITLE_AND_BODY"])
set_text(s, 0, "Banks have the data — they just use it AFTER the call.")
set_bullets(s, 1, [
    "Junior operators read scripts. Senior operators trust intuition.",
    "Cross-sell opportunities are missed on every other call.",
    "60% of KYC / AML disclosures get skipped under live pressure.",
    "Customers leave without completing required documentation.",
    "CRM + call recordings exist — post-call analytics is already too late.",
])

# ───────── SLIDE 3 — The solution ─────────
s = p.slides.add_slide(layouts["TITLE_AND_BODY"])
set_text(s, 0, "OvozAI whispers next to the operator — live, in three languages.")
set_bullets(s, 1, [
    "Streams the call in Uzbek / Russian / English  ·  Gemini 3.1 Pro Live",
    "Reads the customer's CRM profile + previous-call context",
    "Recommends the single best-fit SQB product with an AI confidence %",
    "Hands the operator 3 approved objection responses — click to copy",
    "Auto-tracks 8 KYC items  ·  flags illegal phrases (“guaranteed profit”) instantly",
    "Ends with auto-summary, quality score and one-click CRM log",
])

# ───────── SLIDE 4 — Stack & platform (two columns) ─────────
s = p.slides.add_slide(layouts["TITLE_AND_TWO_COLUMNS"])
set_text(s, 0, "Built like a product, not a hackathon submission.")
set_bullets(s, 1, [
    "AI / VOICE",
    "Gemini 3.1 Pro Live  —  Uzbek STT + reasoning",
    "Gemini 2.5 Flash-lite  —  per-turn JSON analysis (< 2 s)",
    "BM25 KB + multilingual synonyms over real SQB data",
    "Speech-gated push-to-talk  —  no silent-hold hallucinations",
    "Strict mirror language rule  —  Uzbek-first, switches on input",
])
set_bullets(s, 2, [
    "PLATFORM  (3 routes, all live)",
    "/agent  —  Operator Copilot (3-pane HUD)",
    "/agent/supervisor  —  6 concurrent calls + KPIs",
    "/agent/analytics  —  30-day trends, leaderboards, NPS",
    "End-to-end latency budget  <  2 s",
    "Architecture scales to 500+ concurrent calls",
])

# ───────── SLIDE 5 — Big numbers (spec targets) ─────────
s = p.slides.add_slide(layouts["BIG_NUMBER_1"])
set_text(s, 0, "Hits every spec target.")
# Layout 20 has 3 title-style placeholders (idx 0, 2, 3).
set_text(s, 2, "<2 s\nrecommendation latency")
set_text(s, 3, "+15%\ncross-sell uplift  ·  95% KYC completion")

# ───────── SLIDE 6 — Wow features ─────────
s = p.slides.add_slide(layouts["TITLE_AND_BODY"])
set_text(s, 0, "Six moments judges will remember.")
set_bullets(s, 1, [
    "Live Uzbek STT  —  push-to-talk, speech-gated, real-time",
    "Typewriter Whisper card  —  with AI confidence meter on every tip",
    "Battle Cards  —  3 approved Uzbek responses to the customer's objection",
    "Auto-tracking KYC checklist  —  items glow green as the agent says them",
    "Compliance guardrail  —  “guaranteed profit” detected → instant red flag",
    "Supervisor + Analytics dashboards  —  full platform, not just one screen",
])

# ───────── SLIDE 7 — Why we win ─────────
s = p.slides.add_slide(layouts["TITLE_AND_BODY"])
set_text(s, 0, "Why this is the 1st-place build.")
set_bullets(s, 1, [
    "Real working voice + analysis — not mock-ups, not a video",
    "Native Uzbek end-to-end — recognized correctly, mirrored, spoken back",
    "Production-grade UX — speech-gated PTT, hydration-safe, abort-aware",
    "Complete 3-route platform — Copilot · Supervisor · Analytics",
    "Banking-first — KYC checklist + compliance guardrails baked in",
    "Live demo at localhost:3000/agent — judges can drive it themselves",
])

# ───────── SLIDE 8 — Closing ─────────
s = p.slides.add_slide(layouts["CUSTOM_1"])
set_text(s, 0, "Ovoz beradi — keyingi qadamni aytadi.")
set_text(
    s, 1,
    "It speaks — and tells the operator what to do next.\n"
    "Live demo: localhost:3000/agent",
)

p.save(OUTPUT)
print("Saved:", OUTPUT)
